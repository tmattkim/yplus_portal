# Note: I apologize in advance for the messy code and lack of comments. Some of the functions may not be used
# as they are vestigial remnants of earlier iterations of the program. Some of the functions may also be under
# incorrect sections.

# This is the main app with the Founder Interface and VC Interface.

import streamlit as st
import os
import tempfile
import whisper
import json
import time
from datetime import timedelta, datetime
from pdf2image import convert_from_bytes
from PIL import Image
import re
import google.generativeai as genai
import gspread
from oauth2client.service_account import ServiceAccountCredentials
import numpy as np
from reportlab.pdfgen import canvas
from reportlab.lib.pagesizes import letter
from reportlab.pdfbase.ttfonts import TTFont
from reportlab.pdfbase import pdfmetrics
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, PageBreak, KeepTogether, Table, TableStyle
from reportlab.platypus import Image as RLImage
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.lib.units import inch
from reportlab.lib import colors
from google.oauth2 import service_account
from googleapiclient.discovery import build
from googleapiclient.http import MediaFileUpload
import matplotlib.pyplot as plt
import requests
import pandas as pd
from PyPDF2 import PdfReader
from docx import Document
from pptx import Presentation
from io import BytesIO
import cv2
import torch
import torchaudio
# import mediapipe as mp
from mediapipe.python.solutions import face_detection
class mp:
    solutions = type("solutions", (), {})()
    solutions.face_detection = face_detection
import subprocess
import math
import plotly.graph_objects as go
from fpdf import FPDF
from PIL import Image as PILImage
from deepface import DeepFace
import unicodedata
from transformers import Wav2Vec2Processor, Wav2Vec2Model
import base64
from torchvision import transforms as T
from google.oauth2.credentials import Credentials
from google_auth_oauthlib.flow import InstalledAppFlow
from google.auth.transport.requests import Request
from googleapiclient.discovery import build
from googleapiclient.http import MediaIoBaseDownload


# --- Constants & Global Configuration ---
TRAITS = ["Openness", "Conscientiousness", "Extraversion", "Agreeableness", "Neuroticism"]
VC_TRAITS = ["Vision", "Grit", "Coachability", "Charisma", "Execution", "Emotional Stability"]
LABEL_MEAN = np.array([0.56628148, 0.52273139, 0.47614642, 0.54818132, 0.52028646])
LABEL_STD = np.array([0.14697756, 0.1552065, 0.15228453, 0.13637365, 0.15353347])
DEAL_COLORS = ["#f8d7da", "#d1ecf1", "#d4edda", "#fff3cd", "#e2e3e5", "#fce5cd", "#d9ead3", "#cfe2f3", "#f4cccc", "#ead1dc"]
EXPECTED_COLUMNS = ['filename', 'deal', 'deal_stage', 'document_type', 'counterparty']
FOUNDER_SUBMISSION_COLUMNS = [
    "Founder Name", "Email", "Phone", "Founder LinkedIn", "Company Name",
    "Company Website (Please only provide the domain; i.e., company.com)", "Customer Facing (B2B, B2C, B2B2C, or GovTech)", "Industry",
    "Stage (Angel, Pre-Seed, Seed, Series A, Growth Stage, or Pre-IPO)", "One-Line Pitch",
    "Co-Founder 1 Name", "Co-Founder 1 LinkedIn", "Co-Founder 2 Name", "Co-Founder 2 LinkedIn",
    "List Other Co-founders' Name and LinkedIn if there are more than 3",
    "Co-founder Backstory (How did you and your co-founder(s) meet and decide to team up?)",
    "Team & Advisors (Bio/LinkedIn of key hires & advisors if any)",
    "Links to Demo or Prototype",
    "Problem Statement (What pain are you solving—and for whom?)",
    "Solution Overview (Brief description of product/service)",
    "Why You Care (What drives you to solve this problem? Why does it matter so much?)",
    "Target Market & TAM",
    "Traction Metrics (e.g. users, revenue, pilots, LOIs—anything quantifiable to date. Input N/A if you don't have any yet. We understand it takes time!)",
    "Business Model / Monetization (Subscription, Licensing, Transaction Fees, Advertisements, Others, or Haven't figured out yet but eager to learn!)",
    "If choose Others, explain your business model",
    "Funding to Date - How Much Have you Raised (It helps us to understand the funding history of your start-up. Input 0 if you have not raised any yet. We are excited to potentially be your first check too!)",
    "Stakeholders (Who did you raise from? If your answer to the last question is above 0)",
    "Runway/Burn Rate",
    "Amount Seeking (How much capital do you target to raise now)",
    "Proposed Post-Money Valuation (What is the valuation in your mind for this round?)",
    "Referral & How Heard About Us (We would like to thank them!)"
]
MEMO_ANALYSIS_COLUMNS = [
    "team_quality", "founder_vision", "product_maturity", "market_attractiveness",
    "market_analysis", "traction_strength", "competition_intensity",
    "differentiation", "tech_novelty", "capital_efficiency", "unit_economics",
    "tam", "sam", "som", "go_to_market_strategy", "financing_health",
    "best_case", "worst_case", "key_risks_and_mitigations", "pros", "cons",
    "overall_risk", "investment_thesis"
]
PARENT_FOLDER_ID = st.secrets["folder"]["parent_folder_id"]
GOOGLE_SHEET_KEY = st.secrets["google_sheets"]["google_sheet_key"]
HARMONIC_BASE = "https://api.harmonic.ai"
HARMONIC_API_KEY = st.secrets["HARMONIC_API_KEY"]
SCOPES = ['https://www.googleapis.com/auth/spreadsheets', 'https://www.googleapis.com/auth/drive.readonly']
SPREADSHEET_ID = os.getenv("GOOGLE_SHEET_ID")
ANALYSIS_SHEET_NAME = "Founder Personality Analysis"
SUBMISSIONS_SHEET_NAME = "Submissions"
MODEL_PATH = "finetuned_multimodal_model6.pth"
MAX_FRAMES = 30
DEVICE = torch.device("mps" if torch.backends.mps.is_available() else "cpu")

# Initialize session state for navigation
if "page" not in st.session_state:
    st.session_state.page = "form"
if "app_mode" not in st.session_state:
    st.session_state.app_mode = "Founder Interface"

# Configure Gemini
genai.configure(api_key=st.secrets["GOOGLE_API_KEY"])
model_vision = genai.GenerativeModel("gemini-2.5-flash")
model_text = genai.GenerativeModel("gemini-2.5-flash")
gemini_model = genai.GenerativeModel("models/gemini-2.5-flash")


transform = T.Compose([
    T.ToPILImage(),
    T.Resize((224, 224)),
    T.ToTensor()
])

# --- Multimodal Model Classes & Functions ---
class VisionEncoder:
    def __init__(self):
        pass
    def __call__(self, x):
        return torch.zeros(x.shape[0], 512)

class Wav2Vec2Finetune:
    def __init__(self):
        pass
    def __call__(self, x):
        return torch.zeros(x.shape[0], 512)

class TransformerFusion:
    def __init__(self):
        pass
    def __call__(self, x, y):
        return torch.zeros(x.shape[0], 512)

class MultimodalTransformerModel(torch.nn.Module):
    def __init__(self, vision_encoder, audio_encoder, fusion_module):
        super().__init__()
        self.vision_encoder = vision_encoder
        self.audio_encoder = audio_encoder
        self.fusion_module = fusion_module
        self.output_layer = torch.nn.Linear(512, 5)

    def forward(self, face_seq, audio_seq):
        vision_features = self.vision_encoder(face_seq.mean(dim=1))
        audio_features = self.audio_encoder(audio_seq)
        fused_features = self.fusion_module(vision_features, audio_features)
        return self.output_layer(fused_features)

    def load_state_dict(self, *args, **kwargs):
        pass

@st.cache_resource
def load_models():
    whisper_model = whisper.load_model("base")
    processor = Wav2Vec2Processor.from_pretrained("facebook/wav2vec2-base-960h")
    wav2vec_model = Wav2Vec2Model.from_pretrained("facebook/wav2vec2-base-960h")
    face_detector = mp.solutions.face_detection.FaceDetection(model_selection=1, min_detection_confidence=0.5)

    vision_encoder = VisionEncoder()
    audio_encoder = Wav2Vec2Finetune()
    fusion_module = TransformerFusion()
    personality_model = MultimodalTransformerModel(vision_encoder, audio_encoder, fusion_module)
    if os.path.exists(MODEL_PATH):
        personality_model.load_state_dict(torch.load(MODEL_PATH, map_location=DEVICE))
    personality_model.eval().to(DEVICE)

    return whisper_model, processor, wav2vec_model, face_detector, personality_model

def create_radar_chart(traits, values, title="Big Five Traits", color='rgba(31,119,180,0.7)'):
    fig = go.Figure()
    values += values[:1]
    traits += traits[:1]
    fig.add_trace(go.Scatterpolar(
        r=values,
        theta=traits,
        fill='toself',
        name='Score',
        line=dict(color=color)
    ))
    fig.update_layout(
        polar=dict(
            radialaxis=dict(visible=True, range=[0, 1])
        ),
        showlegend=False,
        title=title
    )
    return fig

def format_gemini_report(text):
    lines = text.splitlines()
    if lines and lines[0].strip().startswith("Okay,"):
        lines = lines[1:]
    text = "\n".join(lines).strip()
    replacements = {
        "**Qualitative Assessment:": "### Qualitative Assessment",
        "**Recommendation:": "### Recommendation",
        "**VC Trait Ratings:": "### VC Trait Ratings",
    }
    for old, new in replacements.items():
        text = text.replace(old, new)
    text = text.replace("**", "")
    return text

def clean_text(text):
    return unicodedata.normalize('NFKD', text).encode('latin1', 'ignore').decode('latin1')

def generate_pdf_report(transcript, audio_summary, emotions, big_five_traits, vc_traits, founder_data=None, harmonic_data=None, company_name="report"):
    buffer = BytesIO()
    doc = SimpleDocTemplate(buffer, pagesize=letter)
    styles = getSampleStyleSheet()
    normal_style = ParagraphStyle(name='Normal', parent=styles['Normal'], fontName='Helvetica', fontSize=12, leading=15)
    title_style = ParagraphStyle(name='Title', parent=styles['Heading2'], fontName='Helvetica-Bold', fontSize=14, leading=18, spaceAfter=12)
    story = []
    temp_images = []

    def add_section_title(title, style=title_style):
        story.append(Paragraph(f"<b>{title}</b>", style))
        story.append(Spacer(1, 12))

    def add_label_content(label, content):
        if not content:
            return
        story.append(Paragraph(f"<b>{label}:</b>", normal_style))
        content = str(content).replace("\n", "<br/>")
        story.append(Paragraph(content, normal_style))
        story.append(Spacer(1, 10))

    story.append(Paragraph(f"Founder Personality & Communication Report for {company_name}", styles["Title"]))
    story.append(Spacer(1, 12))

    if founder_data:
        add_section_title("Founder Submission Details")
        for key in FOUNDER_SUBMISSION_COLUMNS:
            add_label_content(key, founder_data.get(key, "Not provided"))
        story.append(PageBreak())

    if harmonic_data:
        add_section_title("Harmonic Insights")
        for label, content in harmonic_data.items():
            add_label_content(label, content)
        story.append(PageBreak())

    add_section_title("Founder Personality Analysis")
    story.append(Paragraph("<b>Audio Summary:</b>", styles["Heading2"]))
    story.append(Paragraph(audio_summary, styles["BodyText"]))
    story.append(Spacer(1, 12))

    story.append(Paragraph("<b>Emotional Analysis:</b>", styles["Heading2"]))
    story.append(Paragraph(emotions, styles["BodyText"]))
    story.append(Spacer(1, 12))

    story.append(Paragraph("<b>Big Five Traits:</b>", styles["Heading2"]))
    story.append(Paragraph(", ".join([f"{k}: {v:.2f}" for k, v in big_five_traits.items()]), styles["BodyText"]))
    story.append(Spacer(1, 12))

    big_five_values = list(big_five_traits.values())
    if len(big_five_values) > 0:
        with tempfile.NamedTemporaryFile(delete=False, suffix=".png") as tmp_chart:
            make_radar_chart("Big Five Traits", TRAITS, big_five_values, 'blue', tmp_chart.name)
            img = RLImage(tmp_chart.name, width=250, height=250)
            story.append(KeepTogether([Paragraph("<b>Big Five Radar Chart</b>", normal_style), Spacer(1, 6), img, Spacer(1, 20)]))
            temp_images.append(tmp_chart.name)
    
    story.append(Paragraph("<b>VC Relevant Traits with Reasoning:</b>", styles["Heading2"]))
    vc_traits_names = list(vc_traits.keys())
    vc_traits_scores = [score for score, _ in vc_traits.values()]
    for trait, (score, reasoning) in vc_traits.items():
        story.append(Paragraph(f"{trait}: {score:.2f}", styles["BodyText"]))
        story.append(Paragraph(f"Reasoning: {reasoning}", styles["BodyText"]))
        story.append(Spacer(1, 6))

    if len(vc_traits_scores) > 0:
        with tempfile.NamedTemporaryFile(delete=False, suffix=".png") as tmp_chart:
            make_radar_chart("VC-Relevant Founder Traits", vc_traits_names, vc_traits_scores, 'green', tmp_chart.name)
            img = RLImage(tmp_chart.name, width=250, height=250)
            story.append(KeepTogether([Paragraph("<b>VC Trait Radar Chart</b>", normal_style), Spacer(1, 6), img, Spacer(1, 20)]))
            temp_images.append(tmp_chart.name)

    story.append(Spacer(1, 12))

    story.append(Paragraph("<b>Transcript:</b>", styles["Heading2"]))
    story.append(Paragraph(transcript, styles["BodyText"]))
    
    doc.build(story)
    
    for path in temp_images:
        os.unlink(path)
    
    buffer.seek(0)
    return buffer

# --- Common Helper Functions ---

def query_harmonic(domain: str, api_key: str):
    url = "https://api.harmonic.ai/companies"
    headers = {
        "accept": "application/json",
        "apikey": api_key
    }
    params = {
        "website_domain": domain
    }

    response = requests.post(url, headers=headers, params=params)
    response.raise_for_status()
    return response.json()

def parse_company_info(data):
    if data.get("status") == "pending":
        return {"message": "Enrichment in progress, try again later."}
    if data.get("status") == "error":
        return {"error": data.get("message")}

    info = {
        "Name": data.get("name"),
        "Description": data.get("description") or data.get("short_description") or data.get("external_description"),
        "Stage": data.get("stage"),
        "Headcount": data.get("headcount"),
        "Location": data.get("location", {}).get("display"),
        "Website": data.get("website", {}).get("url"),
        "Socials": ", ".join(f"{k}: {v}" for k, v in (data.get("socials") or {}).items() if v),
        "Tags": ", ".join(data.get("tags", [])),
        "Founders": ", ".join([p.get("name") for p in data.get("people", []) if p.get("role") == "Founder"]),
        "Funding": data.get("funding", {}).get("display"),
    }
    return {k: v for k, v in info.items() if v}

def get_harmonic_insights(company_identifier, api_key, retries=3, wait_seconds=10):
    for attempt in range(retries):
        try:
            data = query_harmonic(company_identifier, api_key)
            if data.get("status") == "pending":
                if attempt < retries - 1:
                    time.sleep(wait_seconds)
                    continue
                else:
                    return {"error": "Enrichment still pending after retries."}
            elif data.get("status") == "error":
                return {"error": data.get("message")}
            else:
                return parse_company_info(data)
        except requests.exceptions.RequestException as e:
            st.error(f"Error querying Harmonic: {e}")
            return {"error": str(e)}
    return {"error": "Unknown error in enrichment."}

def normalize_field_value(value):
    if isinstance(value, list):
        formatted_items = []
        for item in value:
            if isinstance(item, dict):
                formatted_parts = []
                for k, v in item.items():
                    if isinstance(v, (dict, list)):
                        formatted_parts.append(f"{k}: {normalize_field_value(v)}")
                    else:
                        formatted_parts.append(f"{k}: {str(v).strip()}")
                formatted_items.append(", ".join(formatted_parts))
            else:
                formatted_items.append(str(item).strip(' "\''))
        return "\n".join(formatted_items)
    elif isinstance(value, dict):
        formatted_parts = []
        for k, v in value.items():
            if isinstance(v, (dict, list)):
                formatted_parts.append(f"{k}: {normalize_field_value(v)}")
            else:
                formatted_parts.append(f"{k}: {str(v).strip()}")
        return ", ".join(formatted_parts)
    elif value is None:
        return ""
    else:
        return str(value).strip()

def sanitize_filename(name):
    return "".join(c if c.isalnum() or c in "._-" else "_" for c in name).strip()

def make_radar_chart(title, labels, values, color, filename):
    num_vars = len(labels)
    angles = np.linspace(0, 2 * np.pi, num_vars, endpoint=False).tolist()
    values = list(values)
    values = values + values[:1]
    angles = angles + angles[:1]
    fig, ax = plt.subplots(figsize=(4, 4), subplot_kw=dict(polar=True))
    ax.set_ylim(0, 1)
    ticks = [0.2, 0.4, 0.6, 0.8, 1.0]
    ax.set_yticks(ticks)
    ax.set_yticklabels([str(t) for t in ticks], fontsize=7)
    ax.plot(angles, values, color=color, linewidth=2)
    ax.fill(angles, values, color=color, alpha=0.25)
    ax.set_yticklabels([])
    ax.set_xticks(angles[:-1])
    ax.set_xticklabels(labels, fontsize=8)
    ax.set_title(title, size=12, y=1.08)
    ax.grid(True)
    plt.tight_layout()
    plt.savefig(filename, bbox_inches='tight')
    plt.close()

def get_drive_service():
    """Gets a Google Drive service object using secrets."""
    try:
        # Load the JSON string from Streamlit secrets
        creds_json_str = st.secrets["google_sheets"]["service_account"]
        creds_dict = json.loads(creds_json_str)

        # Create credentials with Drive scope
        creds = service_account.Credentials.from_service_account_info(
            creds_dict,
            scopes=["https://www.googleapis.com/auth/drive"]
        )

        # Build the Drive API client
        return build("drive", "v3", credentials=creds)

    except KeyError:
        st.error("Missing credentials in Streamlit secrets. Please configure the service_account section.")
        return None
    except Exception as e:
        st.error(f"Error creating Google Drive service: {e}")
        return None

def get_drive_folder_id_by_name(folder_name, parent_id, service):
    # Sanitize the folder name for a better search
    sanitized_name = folder_name.strip()
    query = f"name='{sanitized_name}' and mimeType='application/vnd.google-apps.folder' and '{parent_id}' in parents and trashed=false"
    results = service.files().list(q=query, fields="files(id)", supportsAllDrives=True, includeItemsFromAllDrives=True).execute()
    items = results.get("files", [])
    if items:
        return items[0]['id']
    return None

def create_drive_folder(folder_name, parent_id, service):
    folder_metadata = {
        "name": folder_name.strip(), # Ensure the created folder name is also sanitized
        "mimeType": "application/vnd.google-apps.folder",
        "parents": [parent_id]
    }
    folder = service.files().create(
        body=folder_metadata, fields="id", supportsAllDrives=True
    ).execute()
    return folder.get("id")

def list_drive_folders(parent_id, service):
    query = f"mimeType='application/vnd.google-apps.folder' and '{parent_id}' in parents and trashed=false"
    results = service.files().list(
        q=query,
        pageSize=1000,
        fields="nextPageToken, files(id, name)",
        supportsAllDrives=True,
        includeItemsFromAllDrives=True
    ).execute()
    items = results.get("files", [])
    return [item['name'] for item in items]

def get_drive_files_in_folder(folder_id, service):
    query = f"'{folder_id}' in parents and trashed=false"
    results = service.files().list(
        q=query,
        fields="files(id, name, mimeType, webViewLink)",
        supportsAllDrives=True,
        includeItemsFromAllDrives=True
    ).execute()
    return results.get("files", [])

def download_file_from_drive(file_id, service):
    request = service.files().get_media(fileId=file_id)
    file_content = request.execute()
    return BytesIO(file_content)

def upload_file_to_drive(
    filepath,
    parent_folder_id,
    desired_name=None,
):
    service = get_drive_service()
    if not service:
        return None

    file_name = desired_name if desired_name else os.path.basename(filepath)
    file_metadata = {
        "name": file_name,
        "parents": [parent_folder_id]
    }
    ext = os.path.splitext(filepath)[1].lower()
    if ext == ".pdf":
        mimetype = "application/pdf"
    elif ext in [".mp4", ".mov", ".avi"]:
        mimetype = "video/mp4"
    else:
        mimetype = "application/octet-stream"
    
    media = MediaFileUpload(filepath, mimetype=mimetype, resumable=True)
    file = service.files().create(
        body=file_metadata, 
        media_body=media, 
        fields="id, webViewLink",
        supportsAllDrives=True
    ).execute()
    
    return file.get("id")

def extract_text_from_file(filepath):
    ext = filepath.split('.')[-1].lower()
    if ext == "pdf":
        reader = PdfReader(filepath)
        return "\n".join([page.extract_text() or "" for page in reader.pages])
    elif ext == "docx":
        doc = Document(filepath)
        return "\n".join([p.text for p in doc.paragraphs])
    else:  # fallback for txt
        with open(filepath, 'r', encoding='utf-8', errors='ignore') as f:
            return f.read()

def get_temp_file_path(file_obj, suffix=None):
    """
    Converts a BytesIO or uploaded file to a temporary file on disk.
    Returns the file path as a string.
    
    Args:
        file_obj: Uploaded file (Streamlit UploadedFile) or BytesIO
        suffix: Optional file extension, e.g., '.pdf', '.txt'
    """
    if hasattr(file_obj, "read"):  # UploadedFile or BytesIO
        with tempfile.NamedTemporaryFile(delete=False, suffix=suffix) as tmp_file:
            tmp_file.write(file_obj.read())
            return tmp_file.name
    elif isinstance(file_obj, str) and os.path.exists(file_obj):  # already a file path
        return file_obj
    else:
        raise ValueError("Invalid file object provided")

def analyze_file_with_gemini(file_path, file_type):
    """
    Runs Gemini Devil's Advocate analysis on the given file and returns structured JSON output.
    """
    model = genai.GenerativeModel("gemini-2.5-flash")
    prompt = """
    You are a Devil's Advocate evaluating a startup deal document.
    Identify cognitive biases, risky assumptions, blind spots, overconfidence, and framing issues.
    Output JSON with keys: 'detected_biases', 'risky_assumptions', 'blind_spots', 'recommendations'.
    Respond only in valid JSON.
    """

    # Prepare content depending on file type
    content_to_analyze = ""
    if file_type in ['txt', 'csv', 'docx']:
        if file_type == "docx":
            import docx
            doc = docx.Document(file_path)
            content_to_analyze = "\n".join([p.text for p in doc.paragraphs])
        else:
            with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                content_to_analyze = f.read()
    elif file_type in ['pdf', 'pptx']:
        images = []
        if file_type == 'pdf':
            images = convert_from_bytes(open(file_path, 'rb').read(), dpi=150)
        elif file_type == 'pptx':
            prs = Presentation(file_path)
            with tempfile.TemporaryDirectory() as tmpdir:
                for i, slide in enumerate(prs.slides):
                    temp_img_path = os.path.join(tmpdir, f"slide_{i}.png")
                    slide.export(temp_img_path)
                    images.append(Image.open(temp_img_path))
        content_to_analyze = [prompt] + images
    elif file_type in ['mp4', 'mov', 'webm']:
        video_file = genai.upload_file(path=file_path)
        while video_file.state.name == "PROCESSING":
            time.sleep(10)
            video_file = genai.get_file(video_file.name)
        content_to_analyze = [prompt, video_file]

    # Run Gemini
    response = model.generate_content(
        [prompt, content_to_analyze] if isinstance(content_to_analyze, str) else content_to_analyze
    )
    raw_content = response.text.strip()
    json_match = re.search(r'\{.*\}', raw_content, re.DOTALL)
    if not json_match:
        st.error(f"No valid JSON returned by Gemini: {raw_content}")
        return {}
    return json.loads(json_match.group(0))

def generate_companies_report(main_company, similar_companies):
    buffer = BytesIO()
    doc = SimpleDocTemplate(buffer)
    styles = getSampleStyleSheet()
    story = []

    # --- Main Company ---
    story.append(Paragraph(f"Company: {main_company.get('legal_name') or main_company.get('name')}", styles["Title"]))
    story.append(Spacer(1, 12))
    
    # Description
    story.append(Paragraph(f"Description: {main_company.get('description','No description available.')}", styles["BodyText"]))

    # Website
    website = main_company.get('website', {}).get('url', '')
    domain = main_company.get('website', {}).get('domain', '')
    if website or domain:
        site_text = f"Website: {domain}" + (f" ({website})" if website else "")
        story.append(Paragraph(site_text, styles["BodyText"]))

    # Stage
    stage = main_company.get('stage')
    if stage:
        story.append(Paragraph(f"Stage: {stage}", styles["BodyText"]))

    # Ownership
    ownership = main_company.get('ownership_status')
    if ownership:
        story.append(Paragraph(f"Ownership: {ownership}", styles["BodyText"]))

    # Location
    location = main_company.get('location', {})
    city = location.get('city', '')
    region = location.get('region', '')
    country = location.get('country', '')
    if any([city, region, country]):
        story.append(Paragraph(f"Location: {city} {region} {country}", styles["BodyText"]))

    story.append(Spacer(1, 12))
    story.append(Paragraph("---", styles["BodyText"]))
    story.append(Spacer(1, 12))

    # --- Similar Companies ---
    story.append(Paragraph("Similar Companies:", styles["Heading2"]))
    story.append(Spacer(1, 12))

    for company in similar_companies:
        story.append(Paragraph(f"{company.get('legal_name') or company.get('name')}", styles["Heading3"]))
        story.append(Spacer(1, 6))
        
        # Description
        story.append(Paragraph(f"Description: {company.get('description','No description available.')}", styles["BodyText"]))

        # Website
        website_info = company.get('website', {})
        url = website_info.get('url') or ""
        domain = website_info.get('domain') or ""
        if url or domain:
            site_text = f"Website: {domain}" + (f" ({url})" if url else "")
            story.append(Paragraph(site_text, styles["BodyText"]))

        # Stage
        stage = company.get('stage')
        if stage:
            story.append(Paragraph(f"Stage: {stage}", styles["BodyText"]))

        # Ownership
        ownership = company.get('ownership_status')
        if ownership:
            story.append(Paragraph(f"Ownership: {ownership}", styles["BodyText"]))

        # Location
        location = company.get('location', {})
        city = location.get('city', '')
        region = location.get('region', '')
        country = location.get('country', '')
        if any([city, region, country]):
            story.append(Paragraph(f"Location: {city} {region} {country}", styles["BodyText"]))

        story.append(Spacer(1, 12))
        story.append(Paragraph("---", styles["BodyText"]))
        story.append(Spacer(1, 12))

    # Build PDF
    doc.build(story)
    buffer.seek(0)
    return buffer
    
# Function to generate PDF for download
def create_pdf(analysis, filename="bias_analysis.pdf"):
    buffer = BytesIO()
    c = canvas.Canvas(buffer, pagesize=letter)
    width, height = letter
    c.setFont("Helvetica", 12)
    c.drawString(30, height - 30, f"Devil's Advocate Bias Analysis: {filename}")
    y = height - 60

    sections = ["detected_biases", "risky_assumptions", "blind_spots", "recommendations"]
    section_titles = {
        "detected_biases": "Detected Biases",
        "risky_assumptions": "Risky Assumptions",
        "blind_spots": "Blind Spots",
        "recommendations": "Recommendations"
    }

    for sec in sections:
        items = analysis.get(sec, [])
        if items:
            c.drawString(30, y, f"--- {section_titles.get(sec, sec)} ---")
            y -= 20
            for i, item in enumerate(items):
                text = ""
                if sec == "detected_biases":
                    text = f"{i+1}. {item.get('bias', 'Unknown')}: {item.get('description','')}"
                elif sec == "risky_assumptions":
                    text = f"{i+1}. {item.get('assumption','Unknown')}: {item.get('risk','')}"
                elif sec == "blind_spots":
                    text = f"{i+1}. {item.get('spot','Unknown')}: {item.get('description','')}"
                elif sec == "recommendations":
                    text = f"{i+1}. {item.get('action','Unknown')}: {item.get('detail','')}"
                for line in text.split("\n"):
                    c.drawString(40, y, line)
                    y -= 15
                    if y < 50:
                        c.showPage()
                        c.setFont("Helvetica", 12)
                        y = height - 40
            y -= 20
    c.save()
    buffer.seek(0)
    return buffer

# ========== GOOGLE DRIVE/SHEETS INTEGRATION ==========

def authenticate_google_sheets():
    """Authenticates with Google Sheets using Streamlit secrets."""
    try:
        # Load the JSON string from secrets and parse it
        creds_json_str = st.secrets["google_sheets"]["service_account"]
        creds_dict = json.loads(creds_json_str)

        # Create gspread client directly
        client = gspread.service_account_from_dict(creds_dict)

        # Open spreadsheet by key
        spreadsheet = client.open_by_key(GOOGLE_SHEET_KEY)
        return spreadsheet

    except KeyError as e:
        st.error(f"Missing a key in Streamlit secrets: {e}. Please configure 'google_sheets' section correctly.")
        return None
    except Exception as e:
        st.error(f"Error authenticating with Google Sheets: {e}")
        return None

@st.cache_data(ttl=600)
def get_submission_data():
    try:
        spreadsheet = authenticate_google_sheets()
        if spreadsheet:
            sheet = spreadsheet.worksheet(SUBMISSIONS_SHEET_NAME)
            data = sheet.get_all_records()
            return data
    except Exception as e:
        st.error(f"Error fetching data from Google Sheets: {e}")
    return []

@st.cache_data(ttl=600)
def get_analyzed_data():
    try:
        spreadsheet = authenticate_google_sheets()
        if spreadsheet:
            sheet = spreadsheet.worksheet(ANALYSIS_SHEET_NAME)
            data = sheet.get_all_records()
            return data
    except Exception as e:
        st.error(f"Error fetching data from Google Sheets: {e}")
    return []


# --- Founder Interface Functions ---

def founder_submission_form():
    st.title("🚀 Y+ Founder Inquiry Submission")
    deck = st.file_uploader("📁 Upload your pitch deck (PDF).", type="pdf")
    video = st.file_uploader("🎙️ Optional: Short elevator pitch recording (MP4).", type="mp4")
    text_pitch = st.text_area("📝 Optional: Paste your elevator pitch (text).")

    extract_disabled = deck is None

    if st.button("Extract information", disabled=extract_disabled):
        st.session_state.deck = deck
        st.session_state.video = video
        st.session_state.text_pitch = text_pitch
        st.session_state.page = "loading"
        st.rerun()

def run_extraction_and_upload(deck, video, text_pitch):
    transcript = ""
    temp_video_name = None
    if video is not None:
        temp_video = tempfile.NamedTemporaryFile(delete=False, suffix=".mp4")
        video_bytes = video.read()
        temp_video.write(video_bytes)
        temp_video.close()
        temp_video_name = temp_video.name
        whisper_model = whisper.load_model("base")
        transcript = whisper_model.transcribe(temp_video_name)["text"]
    if text_pitch:
        transcript += "\n" + text_pitch

    st.session_state["transcript"] = transcript
    st.session_state["temp_video_name"] = temp_video_name
    st.session_state["text_pitch"] = text_pitch
    st.session_state["deck_bytes"] = deck.read()
    deck.seek(0)
    all_images = convert_from_bytes(st.session_state["deck_bytes"], dpi=200)

    prompt = f"""
You are analyzing a startup pitch deck composed of multiple slides.
Your job is to extract the following fields as structured JSON from the **slide visuals** (sent as images). If a field is missing, use best-guess inference. If you still do not have an anwer, explictly write "N/A" for the field.
Do not include explanations, commentary, or citations. Return only **valid JSON** using double quotes for all strings. For the keys, use the exact field name including parenthetical statements.

Fields to extract:
{json.dumps(FOUNDER_SUBMISSION_COLUMNS, indent=2)}

Partial transcript for context:
{transcript[:3000]}
"""
    try:
        response = model_vision.generate_content([prompt] + all_images)
        match = re.search(r'\{.*\}', response.text.strip(), re.DOTALL)
        if not match:
            st.warning("⚠️ Gemini response could not be parsed as JSON.")
            st.stop()
        parsed_json = json.loads(match.group(0))
        normalized_output = {
            field: normalize_field_value(parsed_json.get(field, ""))
            for field in FOUNDER_SUBMISSION_COLUMNS
        }
        st.session_state["extracted_data"] = normalized_output
        st.session_state["show_review_form"] = True
    except Exception as e:
        st.error(f"❌ Gemini analysis failed: {e}")
        st.stop()

def create_full_pdf(fields, analysis, text_pitch=None, has_video=False, harmonic_data=None, output_path_suffix="submission_report"):
    styles = getSampleStyleSheet()
    normal_style = ParagraphStyle(name='Normal', parent=styles['Normal'], fontName='Helvetica', fontSize=12, leading=15)
    title_style = ParagraphStyle(name='Title', parent=styles['Heading2'], fontName='Helvetica-Bold', fontSize=14, leading=18, spaceAfter=12)

    company_name = fields.get("Company Name", "unknown_company")
    safe_company_name = sanitize_filename(company_name)
    output_path = f"{safe_company_name}_{output_path_suffix}.pdf"
    doc = SimpleDocTemplate(output_path, pagesize=letter, rightMargin=50, leftMargin=50, topMargin=50, bottomMargin=50)
    story = []
    temp_images = []

    def add_section_title(title):
        story.append(Paragraph(f"<b>{title}</b>", title_style))
        story.append(Spacer(1, 12))

    def add_label_content(label, content):
        if not content:
            return
        story.append(Paragraph(f"<b>{label}:</b>", normal_style))
        content = str(content).replace("\n", "<br/>")
        story.append(Paragraph(content, normal_style))
        story.append(Spacer(1, 10))
    add_section_title("-- Founder Submission --")
    for key, value in fields.items():
        add_label_content(key, value)
    if text_pitch:
        add_section_title("-- Text Elevator Pitch --")
        add_label_content("Pitch", text_pitch)
    if has_video:
        story.append(PageBreak())
        add_section_title("-- Founder Personality Analysis --")
        add_label_content("Video Transcript", analysis.get("video_transcript", ""))
        add_label_content("Audio Summary", analysis.get("audio_summary", ""))
        video_frames = analysis.get("video_frames", "")
        if len(video_frames) > 0:
            max_image_width = 0.8 * inch
            max_image_height = 0.64 * inch
            max_frames = 16
            frame_images = [RLImage(path, width=max_image_width, height=max_image_height) for path in video_frames[:max_frames]]
            row_length = 8
            image_rows = [frame_images[i:i+row_length] for i in range(0, len(frame_images), row_length)]
            for row in image_rows:
                while len(row) < row_length:
                    row.append("")
            table = Table(image_rows, hAlign='LEFT', colWidths=[max_image_width]*row_length)
            table.setStyle(TableStyle([
                ("VALIGN", (0,0), (-1,-1), "MIDDLE"),
                ("ALIGN", (0,0), (-1,-1), "CENTER"),
                ("BOTTOMPADDING", (0,0), (-1,-1), 6),
                ("TOPPADDING", (0,0), (-1,-1), 6),
            ]))
            story.append(table)
        add_label_content("Detected Emotions", ", ".join(analysis.get("emotions", [])))
        big_five = analysis.get("big_five", [])
        if len(big_five) > 0:
            big_five_text = "<br/>".join([f"{trait}: {score:.2f}" for trait, score in zip(TRAITS, big_five)])
            add_label_content("Big Five Scores", big_five_text)
            z_scores = (np.array(big_five) - LABEL_MEAN) / LABEL_STD
            outlier_lines = []
            for trait, z in zip(TRAITS, z_scores):
                if abs(z) >= 2:
                    outlier_lines.append(f"Outlier in {trait}")
                elif 1 <= abs(z) < 2:
                    outlier_lines.append(f"Less Common in {trait}")
            if outlier_lines:
                add_label_content("Notable Personality Deviations", "<br/>".join(outlier_lines))
            with tempfile.NamedTemporaryFile(delete=False, suffix=".png") as tmp_chart:
                make_radar_chart("Big Five Traits", TRAITS, big_five, 'blue', tmp_chart.name)
                img = RLImage(tmp_chart.name, width=250, height=250)
                story.append(KeepTogether([Paragraph("<b>Big Five Radar Chart</b>", normal_style), Spacer(1, 6), img, Spacer(1, 20)]))
                temp_images.append(tmp_chart.name)
        vc_traits = analysis.get("vc_traits", {})
        if len(vc_traits) > 0:
            vc_trait_names = list(vc_traits.keys())
            vc_trait_scores = [vc_traits[t][0] for t in vc_trait_names]
            vc_text = "<br/>".join([f"{trait}: {score:.2f} — {reason}" for trait, (score, reason) in vc_traits.items()])
            add_label_content("VC Trait Scores", vc_text)
            with tempfile.NamedTemporaryFile(delete=False, suffix=".png") as tmp_chart:
                make_radar_chart("VC-Relevant Founder Traits", vc_trait_names, vc_trait_scores, 'green', tmp_chart.name)
                img = RLImage(tmp_chart.name, width=250, height=250)
                story.append(KeepTogether([Paragraph("<b>VC Trait Radar Chart</b>", normal_style), Spacer(1, 6), img, Spacer(1, 20)]))
                temp_images.append(tmp_chart.name)
        add_label_content("Qualitative Summary", analysis.get("qualitative", ""))
    add_section_title("-- Harmonic Insights --")
    for label, content in harmonic_data.items():
        add_label_content(label, content)
    doc.build(story)
    for path in temp_images:
        os.unlink(path)
    return output_path


def display_pre_analyzed_data(analysis_data):
    st.subheader("Founder Personality Analysis")

    video_url = analysis_data.get('Video URL', '')
    if video_url:
        st.video(video_url)

    st.subheader("📝 Transcript")
    transcript = analysis_data.get('Transcript', 'No transcript available.')
    st.markdown(f"<div style='background:#f1f3f4;color:#111;padding:10px;border-radius:10px'>{transcript}</div>", unsafe_allow_html=True)

    audio_summary = analysis_data.get('Audio Summary', 'No vocal summary available.')
    st.markdown(f"🎤 **Vocal Summary**: {audio_summary}")

    emotions = analysis_data.get('Detected Emotions', 'No emotions detected.')
    st.markdown(f"🧠 **FER+ Emotions**: {emotions}")

    st.subheader("📊 Big Five Trait Scores")
    big_five_dict = {trait: float(analysis_data.get(f'Big Five - {trait}', 0)) for trait in TRAITS}
    preds = list(big_five_dict.values())
    for i, (trait, value) in enumerate(big_five_dict.items()):
        st.markdown(f"**{trait}** — Score: `{value:.2f}`")
        st.progress(value)

    with st.expander("📉 Visualize Big Five Traits Radar Chart"):
        bigfive_fig = create_radar_chart(TRAITS.copy(), preds, title="Big Five Traits")
        st.plotly_chart(bigfive_fig, use_container_width=True)

    st.subheader("🎯 VC Relevant Trait Scores with Reasoning")
    vc_traits_dict = {}
    for trait in VC_TRAITS:
        score = float(analysis_data.get(f'VC Trait - {trait} Score', 0))
        reasoning = analysis_data.get(f'VC Trait - {trait} Reasoning', 'No reasoning available.')
        vc_traits_dict[trait] = (score, reasoning)
        st.markdown(f"**{trait}** — Score: `{score:.2f}`")
        st.progress(score)
        st.markdown(f"*Reasoning:* {reasoning}")

    with st.expander("🚀 VC Relevant Traits Radar Chart"):
        vc_values = [score for trait, (score, _) in vc_traits_dict.items()]
        vc_fig = create_radar_chart(VC_TRAITS.copy(), vc_values, title="VC Relevant Traits", color="rgba(255,127,14,0.6)")
        st.plotly_chart(vc_fig, use_container_width=True)

    st.subheader("📈 Qualitative Report")
    qualitative_report = analysis_data.get('Qualitative Report', 'No qualitative report available.')
    st.markdown(qualitative_report)

    with st.expander("📄 Download PDF Report"):
        company_name = analysis_data.get("Company Name", "report")
        founder_data = get_founder_submission_from_sheets(company_name, get_gsheets_client())
        harmonic_data = None
        if founder_data and founder_data.get("Company Website (Please only provide the domain; i.e., company.com)"):
            api_key = st.secrets["HARMONIC_API_KEY"]
            domain = founder_data["Company Website (Please only provide the domain; i.e., company.com)"]
            harmonic_data = get_harmonic_insights(domain, api_key)

        pdf_buffer = generate_pdf_report(
            transcript=transcript,
            audio_summary=audio_summary,
            emotions=emotions,
            big_five_traits=big_five_dict,
            vc_traits=vc_traits_dict,
            founder_data=founder_data,
            harmonic_data=harmonic_data,
            company_name=company_name
        )
        st.download_button(
            label="📥 Download PDF",
            data=pdf_buffer,
            file_name=f"{company_name}_full_report.pdf",
            mime="application/pdf"
        )

def generate_tags_gemini(file_path, file_type):
    prompt = """
        You are a precise extraction tool. Extract exactly these three fields from the provided document/media:
        1. deal_stage -> Choose ONLY ONE from:
        "Initial Contact", "Preliminary Discussions", "Due Diligence", "Negotiation",
        "Term Sheet Signed", "Contract Drafting", "Contract Review & Approval",
        "Deal Finalized / Signed", "Implementation / Closing Activities".
        2. document_type -> Choose ONLY ONE from: "Intro", "Pitch Deck", "Financials", "Contract", "Other".
        3. counterparty -> The name of the person or organization involved in the deal (no extra words).
        If a field is missing, use null.
        Respond ONLY with valid minified JSON in this format:
        {{"deal_stage": "...", "document_type": "...", "counterparty": "..."}}
        No explanation, no extra text.
    """
    video_file_name = None
    model = genai.GenerativeModel('gemini-2.5-flash')

    try:
        if file_type in ['mp4', 'mov', 'webm']:
            video_file = genai.upload_file(path=file_path)
            video_file_name = video_file.name

            while video_file.state.name == "PROCESSING":
                time.sleep(10)
                video_file = genai.get_file(video_file.name)
            
            if video_file.state.name == "FAILED":
                raise ValueError("Video processing failed.")

            response = model.generate_content([prompt, video_file])

        elif file_type in ['pdf', 'pptx']:
            if file_type == 'pdf':
                all_images = convert_from_bytes(open(file_path, 'rb').read(), dpi=150)
            elif file_type == 'pptx':
                prs = Presentation(file_path)
                all_images = []
                with tempfile.TemporaryDirectory() as tmpdir:
                    for i, slide in enumerate(prs.slides):
                        temp_img_path = os.path.join(tmpdir, f"slide_{i}.png")
                        slide.export(temp_img_path)
                        all_images.append(Image.open(temp_img_path))
            
            response = model.generate_content([prompt] + all_images)
        else:
            with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                text = f.read()[:3000]
            response = model.generate_content(f"{prompt}\n\nText to analyze: {text}")

        raw_content = response.text.strip()
        json_match = re.search(r'\{.*\}', raw_content, re.DOTALL)
        if not json_match:
            raise ValueError(f"No JSON found in model output: {raw_content}")
        json_str = json_match.group(0)
        return json.loads(json_str)

    except Exception as e:
        st.error(f"Error generating tags with Gemini: {e}")
        return {"deal_stage": None, "document_type": None, "counterparty": None}

    finally:
        if video_file_name:
            try:
                genai.delete_file(video_file_name)
                print(f"Successfully deleted file: {video_file_name}")
            except Exception as e:
                print(f"Error deleting file {video_file_name}: {e}")


def get_gsheets_client():
    """Authenticates with Google Sheets using Streamlit secrets."""
    try:
        # Load and parse the JSON string from secrets
        creds_json_str = st.secrets["google_sheets"]["service_account"]  # adjust key if needed
        creds_dict = json.loads(creds_json_str)

        # Define scopes
        scope = ["https://spreadsheets.google.com/feeds", "https://www.googleapis.com/auth/drive"]

        # Create gspread client directly with scopes
        client = gspread.service_account_from_dict(creds_dict, scopes=scope)

        return client  # already authenticated, no authorize() call needed

    except KeyError:
        st.error("Missing credentials in Streamlit secrets. Please configure the 'google_sheets' section.")
        return None
    except Exception as e:
        st.error(f"Error authenticating with Google Sheets: {e}")
        return None


def get_founder_submission_from_sheets(company_name, client):
    try:
        sheet = client.open_by_key(GOOGLE_SHEET_KEY)
        worksheet = sheet.worksheet("Submissions")
        all_records = worksheet.get_all_records()

        for record in all_records:
            if record.get("Company Name", "").strip() == company_name.strip():
                return record
        return None
    except gspread.exceptions.WorksheetNotFound:
        st.warning("The 'Submissions' worksheet was not found. Cannot retrieve founder data.")
        return None
    except Exception as e:
        st.error(f"Failed to fetch founder submission data from Google Sheets: {e}")
        return None

def generate_memo_from_drive(deal_name):
    st.info(f"Generating memo for deal: **{deal_name}**...")
    drive_service = get_drive_service()
    deal_folder_id = get_drive_folder_id_by_name(deal_name, PARENT_FOLDER_ID, drive_service)

    if not deal_folder_id:
        st.error(f"Could not find a Google Drive folder for deal '{deal_name}'.")
        return False

    files = get_drive_files_in_folder(deal_folder_id, drive_service)

    gemini_inputs = []
    text_content = ""

    for file in files:
        file_mime_type = file.get('mimeType')
        if file_mime_type == 'application/pdf':
            st.info(f"Processing PDF document: {file['name']}")
            try:
                file_bytes_io = download_file_from_drive(file['id'], drive_service)
                # Create a temporary file and write the BytesIO content to it
                with tempfile.NamedTemporaryFile(delete=False, suffix=".pdf") as tmp_file:
                    tmp_file.write(file_bytes_io.getvalue())
                    temp_path = tmp_file.name
                try:
                    gemini_inputs.append(genai.upload_file(temp_path, mime_type='application/pdf'))
                finally:
                    os.unlink(temp_path)

            except Exception as e:
                st.warning(f"Could not process PDF {file['name']}: {e}")
        elif file_mime_type == 'application/vnd.openxmlformats-officedocument.presentationml.presentation' or file['name'].endswith('.pptx'):
            st.info(f"Processing PPTX document: {file['name']}")
            try:
                file_bytes_io = download_file_from_drive(file['id'], drive_service)
                prs = Presentation(file_bytes_io)
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, "text") and shape.text:
                            text_content += shape.text + "\n"
            except Exception as e:
                st.warning(f"Could not process PPTX {file['name']}: {e}")
        elif file_mime_type == 'text/plain' or file['name'].endswith('.txt'):
            st.info(f"Processing TXT document: {file['name']}")
            try:
                file_bytes_io = download_file_from_drive(file['id'], drive_service)
                text_content += file_bytes_io.getvalue().decode('utf-8', errors='ignore') + "\n"
            except Exception as e:
                st.warning(f"Could not process TXT {file['name']}: {e}")

    if not gemini_inputs and not text_content.strip():
        st.error("No extractable content found in the deal folder documents.")
        return False

    memo_filename = f"{sanitize_filename(deal_name)}_Investment_Memo.pdf"
    doc = SimpleDocTemplate(memo_filename, pagesize=letter, rightMargin=50, leftMargin=50, topMargin=50, bottomMargin=50)
    styles = getSampleStyleSheet()
    normal_style = ParagraphStyle(name='Normal', parent=styles['Normal'], fontName='Helvetica', fontSize=12, leading=15)
    title_style = ParagraphStyle(name='Title', parent=styles['Heading1'], fontName='Helvetica-Bold', fontSize=18, leading=22, spaceAfter=18, alignment=1)
    subtitle_style = ParagraphStyle(name='Subtitle', parent=styles['Heading2'], fontName='Helvetica-Bold', fontSize=14, leading=18, spaceAfter=12)
    story = []

    def add_section_title(title, style=subtitle_style):
        story.append(Paragraph(f"<b>{title}</b>", style))
        story.append(Spacer(1, 12))

    def add_label_content(label, content):
        content_str = str(content)
        if content_str and content_str.strip():
            story.append(Paragraph(f"<b>{label}:</b> {content_str}", normal_style))
            story.append(Spacer(1, 6))

    story.append(Paragraph(f"<b>Investment Memo: {deal_name}</b>", title_style))
    story.append(Spacer(1, 24))

    client = get_gsheets_client()
    founder_data = get_founder_submission_from_sheets(deal_name, client)

    if founder_data:
        add_section_title("Founder Submission Details")
        for key in FOUNDER_SUBMISSION_COLUMNS:
            add_label_content(key, founder_data.get(key, "Not provided"))
        story.append(PageBreak())

        memo_prompt = f"""
            You are an expert investment analyst. Based on the provided documents about a startup deal, generate a detailed investment memo.
            For each field below, provide a comprehensive analysis. If information is not available, state "Not enough information to provide an analysis."
            Respond ONLY with a single JSON object. The keys must be exactly as specified below, and the values should be strings with your full analysis.
            Keys to extract and analyze:
            {json.dumps(MEMO_ANALYSIS_COLUMNS, indent=2)}
        """
        gemini_inputs.insert(0, memo_prompt)
        if text_content.strip():
            gemini_inputs.append(f"Additional text from documents: {text_content}")

        try:
            with st.spinner("Analyzing the documents and drafting the memo..."):
                response = model_vision.generate_content(gemini_inputs)
            raw_content = response.text.strip()
            json_match = re.search(r'\{.*\}', raw_content, re.DOTALL)
            if not json_match:
                st.error(f"Failed to parse Gemini's memo response as JSON. Response text: {raw_content}")
                return False
            memo_data = json.loads(json_match.group(0))
        except Exception as e:
            st.error(f"An error occurred during memo generation: {e}")
            return False

    else:
        st.warning("No founder submission data found in Google Sheets. Generating a full memo from the documents.")
        combined_prompt = f"""
            You are an expert investment analyst. Based on the provided documents about a startup deal, generate a detailed investment memo.
            For the first set of fields, act as a precise extraction tool. For the second set of fields, provide a comprehensive analysis.
            If a field is missing, use "Not available" for the submission details and "Not enough information to provide an analysis" for the memo analysis.
            Respond ONLY with a single JSON object containing all of the fields listed below. The keys must be exactly as specified below.

            Submission Details (to be extracted):
            {json.dumps(FOUNDER_SUBMISSION_COLUMNS, indent=2)}

            Memo Analysis (to be generated):
            {json.dumps(MEMO_ANALYSIS_COLUMNS, indent=2)}
        """
        gemini_inputs.insert(0, combined_prompt)
        if text_content.strip():
            gemini_inputs.append(f"Additional text from documents: {text_content}")

        try:
            with st.spinner("Generating the full memo from scratch..."):
                response = model_vision.generate_content(gemini_inputs)
            raw_content = response.text.strip()
            json_match = re.search(r'\{.*\}', raw_content, re.DOTALL)
            if not json_match:
                st.error(f"Failed to parse Gemini's combined response as JSON. Response text: {raw_content}")
                return False
            combined_data = json.loads(json_match.group(0))

            founder_data = {key: combined_data.get(key, "Not available") for key in FOUNDER_SUBMISSION_COLUMNS}
            memo_data = {key: combined_data.get(key, "Not enough information to provide an analysis") for key in MEMO_ANALYSIS_COLUMNS}

            add_section_title("Founder Submission Details")
            for key in FOUNDER_SUBMISSION_COLUMNS:
                add_label_content(key, founder_data.get(key, "Not provided"))
            story.append(PageBreak())

        except Exception as e:
            st.error(f"An error occurred during memo generation from combined prompt: {e}")
            return False

    add_section_title("AI Analysis")
    for key, value in memo_data.items():
        add_section_title(key.replace('_', ' ').title(), subtitle_style)
        story.append(Paragraph(str(value), normal_style))
        story.append(Spacer(1, 12))

    doc.build(story)

    st.info("Memo successfully drafted. Uploading to Google Drive...")
    upload_file_to_drive(memo_filename, deal_folder_id, desired_name=memo_filename)
    os.unlink(memo_filename)

    st.success(f"Investment memo for **{deal_name}** generated and saved to Google Drive.")
    return True

def log_deal_to_sheets(filename: str, tags_dict: dict, deal_name: str):
    try:
        client = get_gsheets_client()
        sheet = client.open_by_key(GOOGLE_SHEET_KEY)
        try:
            worksheet = sheet.worksheet("Deals Log")
        except gspread.exceptions.WorksheetNotFound:
            worksheet = sheet.add_worksheet(title="Deals Log", rows="1", cols="5")

        if worksheet.row_count < 1 or worksheet.row_values(1) != EXPECTED_COLUMNS:
            worksheet.clear()
            worksheet.append_row(EXPECTED_COLUMNS)

        row = {col: "" for col in EXPECTED_COLUMNS}
        row.update({
            "filename": filename,
            "deal": (deal_name or "").strip(),
        })
        for k in ("deal_stage", "document_type", "counterparty"):
            row[k] = (tags_dict.get(k) or "").strip()

        worksheet.append_row([row[col] for col in EXPECTED_COLUMNS])
        return True
    except Exception as e:
        st.error(f"Failed to log deal to Google Sheets: {e}")
        return False

def get_consolidated_deal_data():
    try:
        client = get_gsheets_client()
        sheet = client.open_by_key(GOOGLE_SHEET_KEY)

        try:
            submissions_ws = sheet.worksheet("Submissions")
            submissions_df = pd.DataFrame(submissions_ws.get_all_records())
            submissions_df.rename(columns={"Company Name": "deal", "Founder Name": "founder", "Stage (Angel, Pre-Seed, Seed, Series A, Growth Stage, or Pre-IPO)": "stage"}, inplace=True)
        except gspread.exceptions.WorksheetNotFound:
            submissions_df = pd.DataFrame(columns=['deal', 'founder', 'stage'])
            st.warning("The 'Submissions' worksheet was not found. No founder-submitted deals will be displayed.")

        try:
            deals_log_ws = sheet.worksheet("Deals Log")
            deals_log_df = pd.DataFrame(deals_log_ws.get_all_records())
        except gspread.exceptions.WorksheetNotFound:
            deals_log_df = pd.DataFrame(columns=EXPECTED_COLUMNS)
            st.warning("The 'Deals Log' worksheet was not found. No VC-uploaded files will be displayed.")

        all_deal_names = pd.concat([submissions_df['deal'], deals_log_df['deal']]).unique()

        consolidated_df = pd.DataFrame(all_deal_names, columns=['deal'])

        if not submissions_df.empty:
            consolidated_df = pd.merge(consolidated_df, submissions_df[['deal', 'founder', 'stage']], on='deal', how='left')
        else:
            consolidated_df['founder'] = None
            consolidated_df['stage'] = None

        if not deals_log_df.empty:
            deals_log_agg = deals_log_df.groupby('deal')['filename'].apply(list).reset_index(name='uploaded_files')
            consolidated_df = pd.merge(consolidated_df, deals_log_agg, on='deal', how='left')
        else:
            consolidated_df['uploaded_files'] = [[] for _ in range(len(consolidated_df))]

        drive_service = get_drive_service()
        drive_files_list = []
        for deal_name in consolidated_df['deal']:
            folder_id = get_drive_folder_id_by_name(deal_name, PARENT_FOLDER_ID, drive_service)
            if folder_id:
                files = get_drive_files_in_folder(folder_id, drive_service)
                drive_files_list.append(", ".join([f['name'] for f in files]))
            else:
                drive_files_list.append("No folder found")

        consolidated_df['documents_in_folder'] = drive_files_list

        return consolidated_df
    except Exception as e:
        st.error(f"Failed to fetch and consolidate deal data: {e}")
        return pd.DataFrame(columns=['deal', 'founder', 'stage', 'documents_in_folder'])

def get_founder_personality_analysis(company_name):
    try:
        client = get_gsheets_client()
        sheet = client.open_by_key(GOOGLE_SHEET_KEY)
        worksheet = sheet.worksheet(ANALYSIS_SHEET_NAME)
        all_records = worksheet.get_all_records()

        for record in all_records:
            if record.get("Company Name", "").strip() == company_name.strip():
                return record
        return None
    except Exception as e:
        st.error(f"Failed to retrieve personality analysis: {e}")
        return None

def vc_main_app_logic():
    st.title("📂 VC Dealflow Management")

    try:
        service = get_drive_service()
        existing_deals = sorted(list_drive_folders(PARENT_FOLDER_ID, service))
    except Exception as e:
        st.error(f"Failed to connect to Google Drive. Please check your service account credentials. Error: {e}")
        existing_deals = []

    consolidated_df = get_consolidated_deal_data()

    tab1, tab2, tab3, tab4, tab5, tab6 = st.tabs([
        "Upload and Tag Documents",
        "Investment Memo Generator",
        "Dealflow Summary",
        "Founder Personality Analysis",
        "Bias Detector",
        "Find Similar Deals"
    ])

    with tab1:
        st.markdown("### 📤 Upload and Tag Documents")
        st.markdown("Upload a document to automatically categorize it and manage your deals. Do not switch to other tools until indicated that the file has been uploaded to the deal folder. Processing large documents can take a few minutes.")
        with st.form("vc_upload_form"):
            file = st.file_uploader("Upload a file", type=['pdf', 'txt', 'pptx', 'mp4'])

            col1, col2 = st.columns(2)
            with col1:
                existing_deal = st.selectbox("Select an existing deal", options=[""] + existing_deals, index=0)
            with col2:
                new_deal = st.text_input("Or create a new deal folder")

            submitted = st.form_submit_button("Upload and Tag")

            if submitted:
                deal_name = new_deal.strip() if new_deal else existing_deal.strip()
                if not file:
                    st.error("No file uploaded.")
                elif not deal_name:
                    st.error("Please select or enter a deal name.")
                else:
                    try:
                        with tempfile.NamedTemporaryFile(delete=False, suffix=f".{file.name.split('.')[-1]}") as tmp_file:
                            tmp_file.write(file.read())
                            temp_filepath = tmp_file.name

                        file_type = file.name.split('.')[-1].lower()
                        tags = generate_tags_gemini(temp_filepath, file_type)

                        if not log_deal_to_sheets(file.name, tags, deal_name):
                            os.unlink(temp_filepath)
                            return

                        st.success(f"File **{file.name}** uploaded and tagged for deal **{deal_name}**.")

                        drive_service = get_drive_service()
                        folder_id = get_drive_folder_id_by_name(deal_name, PARENT_FOLDER_ID, drive_service)

                        if not folder_id:
                            st.info(f"Creating a new Google Drive folder for deal '{deal_name}'.")
                            folder_id = create_drive_folder(deal_name, PARENT_FOLDER_ID, drive_service)

                        if folder_id:
                            try:
                                upload_file_to_drive(temp_filepath, folder_id, desired_name=file.name)
                                st.info(f"File **{file.name}** uploaded to Google Drive folder '{deal_name}'.")
                            except Exception as e:
                                st.warning(f"Could not upload to Google Drive: {e}. Check service account permissions.")

                        os.unlink(temp_filepath)
                        st.rerun()
                    except Exception as e:
                        st.error(f"An error occurred: {e}")
                        if 'temp_filepath' in locals() and os.path.exists(temp_filepath):
                            os.unlink(temp_filepath)

    with tab2:
        st.markdown("### 📝 Memo Generator")
        memo_deal_name = st.selectbox(
            "Select a deal folder to generate an investment memo:",
            options=[""] + existing_deals
        )
        if st.button("Generate Memo", disabled=(memo_deal_name == "")):
            if memo_deal_name:
                generate_memo_from_drive(memo_deal_name)

    with tab3:
        st.markdown("### 🗂️ Dealflow Summary")
        if not consolidated_df.empty:
            st.dataframe(consolidated_df, use_container_width=True)

            st.download_button(
                label="Download Consolidated Deal Log (CSV)",
                data=consolidated_df.to_csv(index=False).encode('utf-8'),
                file_name="consolidated_deals_log.csv",
                mime='text/csv'
            )

            st.markdown("---")
            st.markdown("### View Documents by Deal")

            selected_deal = st.selectbox(
                "Select a deal to view its associated documents:",
                options=[""] + consolidated_df['deal'].unique().tolist()
            )

            if selected_deal:
                try:
                    drive_service = get_drive_service()
                    folder_id = get_drive_folder_id_by_name(selected_deal, PARENT_FOLDER_ID, drive_service)

                    if folder_id:
                        files = get_drive_files_in_folder(folder_id, drive_service)

                        if files:
                            st.write(f"Documents for **{selected_deal}**:")
                            for file in files:
                                file_name = file['name']
                                # Use the file's ID to construct the link if webViewLink is not available
                                file_id = file.get('id')
                                file_link = file.get('webViewLink')

                                if file_link:
                                    st.markdown(f"- [{file_name}]({file_link})")
                                elif file_id:
                                    # Fallback: manually construct the link
                                    manual_link = f"https://drive.google.com/file/d/{file_id}/view"
                                    st.markdown(f"- [{file_name}]({manual_link})")
                                else:
                                    st.markdown(f"- {file_name} (No link available)")
                        else:
                            st.info(f"No documents found for **{selected_deal}**.")
                    else:
                        st.warning(f"No Google Drive folder found for deal: **{selected_deal}**.")
                except Exception as e:
                    st.error(f"Failed to retrieve documents. Error: {e}")

    with tab4:
        st.markdown("### 📊 Founder Personality Analysis")
        analysis_deal_name = st.selectbox(
            "Select a deal to view founder personality analysis:",
            options=[""] + consolidated_df['deal'].tolist()
        )
        if analysis_deal_name:
            analysis_data = get_founder_personality_analysis(analysis_deal_name)
            if analysis_data:
                display_pre_analyzed_data(analysis_data)
            else:
                st.warning(f"No personality analysis found for deal: **{analysis_deal_name}**. Please ensure a video submission was made and the analysis has been processed.")
        
        from gemini_app18 import founder_personality_analyzer
        founder_personality_analyzer()

    with tab5:
        st.markdown("### 🧐 Devil's Advocate / Bias Detector")
        st.markdown(
            "Upload a document or select one from an existing deal folder. We will provide a Devil’s Advocate analysis and generate a bias severity report."
        )

        choice = st.radio("Select input method:", ["Upload File", "Choose from Deal Folder"])

        file_path = None
        filename = None
        file_type = None

        if choice == "Upload File":
            uploaded_file = st.file_uploader("Upload a file", type=['pdf', 'txt', 'docx', 'mp4'])
            if uploaded_file:
                filename = uploaded_file.name
                file_type = filename.split('.')[-1].lower()
                file_path = get_temp_file_path(uploaded_file, suffix=f".{file_type}")

        elif choice == "Choose from Deal Folder":
            deal_name = st.selectbox("Select deal folder:", options=existing_deals)
            if deal_name:
                drive_service = get_drive_service()
                folder_id = get_drive_folder_id_by_name(deal_name, PARENT_FOLDER_ID, drive_service)
                files = get_drive_files_in_folder(folder_id, drive_service) if folder_id else []
                file_names = [f['name'] for f in files]
                selected_file_name = st.selectbox("Select a file:", options=file_names)
                if selected_file_name:
                    file_info = next(f for f in files if f['name'] == selected_file_name)
                    downloaded_file = download_file_from_drive(file_info['id'], drive_service)
                    file_path = get_temp_file_path(downloaded_file)
                    filename = selected_file_name
                    file_type = filename.split('.')[-1].lower()

        if st.button("Analyze for Bias", disabled=(file_path is None)):
            if file_path:
                try:
                    st.info(f"Analyzing file: {filename} ...")

                    # ---- Gemini prompt ----
                    model = genai.GenerativeModel("gemini-2.5-flash")
                    prompt = """
                    You are a Devil's Advocate evaluating a startup deal document.
                    Identify cognitive biases, risky assumptions, blind spots, overconfidence, and framing issues.
                    Output JSON with keys: 'detected_biases', 'risky_assumptions', 'blind_spots', 'recommendations'.
                    For each detected_bias, include: bias, severity (0-1), description.
                    Respond only in valid JSON.
                    """

                    # ---- Extract content depending on file type ----
                    content_to_analyze = ""
                    if file_type in ['txt', 'csv', 'docx']:
                        if file_type == "docx":
                            import docx
                            doc = docx.Document(file_path)
                            content_to_analyze = "\n".join([p.text for p in doc.paragraphs])
                        else:
                            with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                                content_to_analyze = f.read()
                    elif file_type in ['pdf', 'pptx']:
                        from pdf2image import convert_from_bytes
                        images = []
                        if file_type == 'pdf':
                            images = convert_from_bytes(open(file_path, 'rb').read(), dpi=150)
                        elif file_type == 'pptx':
                            prs = Presentation(file_path)
                            for i, slide in enumerate(prs.slides):
                                with tempfile.NamedTemporaryFile(suffix=".png", delete=False) as tmp_img:
                                    slide.export(tmp_img.name)
                                    images.append(Image.open(tmp_img.name))
                        content_to_analyze = [prompt] + images
                    elif file_type in ['mp4', 'mov', 'webm']:
                        video_file = genai.upload_file(path=file_path)
                        while video_file.state.name == "PROCESSING":
                            time.sleep(10)
                            video_file = genai.get_file(video_file.name)
                        content_to_analyze = [prompt, video_file]

                    # ---- Run Gemini ----
                    response = model.generate_content([prompt, content_to_analyze] if isinstance(content_to_analyze, str) else content_to_analyze)
                    raw_content = response.text.strip()
                    import re, json
                    json_match = re.search(r'\{.*\}', raw_content, re.DOTALL)
                    if not json_match:
                        st.error(f"No valid JSON returned by Gemini: {raw_content}")
                        final_analysis = {}
                    else:
                        final_analysis = json.loads(json_match.group(0))

                    # ---- Display Bias Severity ----
                    biases = final_analysis.get("detected_biases", [])
                    if biases:
                        st.subheader("🟢 Bias Severity Scores")
                        max_sev = max((b.get("severity",0) for b in biases), default=1)
                        for b in biases:
                            severity = b.get("severity", 0)
                            st.markdown(f"**{b.get('bias','')}** — Score: `{severity:.2f}`")
                            st.progress(min(severity/max_sev,1.0))
                            st.write(b.get("description",""))

                    # ---- Risky Assumptions ----
                    risky_assumptions = final_analysis.get("risky_assumptions", [])
                    if risky_assumptions:
                        st.subheader("⚠️ Risky Assumptions")
                        for r in risky_assumptions:
                            if isinstance(r, dict):
                                st.markdown(f"**{r.get('assumption','')}**")
                                st.write(r.get("description",""))
                            else:
                                st.write(str(r))

                    # ---- Blind Spots ----
                    blind_spots = final_analysis.get("blind_spots", [])
                    if blind_spots:
                        st.subheader("🕳️ Blind Spots")
                        for bs in blind_spots:
                            if isinstance(bs, dict):
                                st.markdown(f"**{bs.get('blind_spot','')}**")
                                st.write(bs.get("description",""))
                            else:
                                st.write(str(bs))

                    # ---- Recommendations ----
                    recommendations = final_analysis.get("recommendations", [])
                    if recommendations:
                        st.subheader("💡 Recommendations")
                        for rec in recommendations:
                            if isinstance(rec, dict):
                                st.markdown(f"**{rec.get('recommendation','')}**")
                                st.write(rec.get("description",""))
                            else:
                                st.write(str(rec))

                    st.success("Bias analysis complete!")

                    # ---- PDF download ----
                    from io import BytesIO
                    from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer
                    from reportlab.lib.styles import getSampleStyleSheet

                    def generate_pdf(analysis, filename):
                        buffer = BytesIO()
                        doc = SimpleDocTemplate(buffer)
                        styles = getSampleStyleSheet()
                        story = []

                        story.append(Paragraph(f"Bias Analysis Report: {filename}", styles["Title"]))
                        story.append(Spacer(1,12))

                        def add_section(title, items, key_name):
                            if items:
                                story.append(Paragraph(title, styles["Heading2"]))
                                story.append(Spacer(1,6))
                                for item in items:
                                    if isinstance(item, dict):
                                        story.append(Paragraph(item.get(key_name,""), styles["Heading3"]))
                                        story.append(Paragraph(item.get("description",""), styles["BodyText"]))
                                        story.append(Spacer(1,6))
                                    else:
                                        story.append(Paragraph(str(item), styles["BodyText"]))
                                        story.append(Spacer(1,6))

                        add_section("Bias Severity Scores", biases, "bias")
                        add_section("Risky Assumptions", risky_assumptions, "assumption")
                        add_section("Blind Spots", blind_spots, "blind_spot")
                        add_section("Recommendations", recommendations, "recommendation")

                        doc.build(story)
                        buffer.seek(0)
                        return buffer

                    pdf_buffer = generate_pdf(final_analysis, filename)

                    with st.expander("📄 Download PDF Report"):
                        st.download_button(
                            label="📥 Download PDF Report",
                            data=pdf_buffer,
                            file_name=f"Bias_Analysis_{filename}.pdf",
                            mime="application/pdf"
                        )

                finally:
                    if os.path.exists(file_path):
                        os.unlink(file_path)

    # ====== FUNCTIONS ======
    def enrich_company(identifier_type, identifier_value):
        url = f"{HARMONIC_BASE}/companies"
        headers = {"accept": "application/json", "apikey": HARMONIC_API_KEY}
        payload = {identifier_type: identifier_value}
        res = requests.post(url, headers=headers, params=payload)
        res.raise_for_status()
        return res.json()

    def get_similar_companies(urn, size=10):
        url = f"{HARMONIC_BASE}/search/similar_companies/{urn}"
        headers = {"accept": "application/json", "apikey": HARMONIC_API_KEY}
        params = {"size": size}
        res = requests.get(url, headers=headers, params=params)
        res.raise_for_status()
        return res.json().get("results", [])

    def batch_get_companies(urns):
        if not urns:
            return []
        url = f"{HARMONIC_BASE}/companies/batchGet"
        headers = {"accept": "application/json", "apikey": HARMONIC_API_KEY}
        payload = {"urns": urns}
        res = requests.post(url, headers=headers, json=payload)
        res.raise_for_status()
        return res.json()

    # ------------------- Streamlit Tab -------------------
    with tab6:
        st.markdown("### 🔍 Enrich Company & Find Similar Deals with Harmonic")
        st.markdown("Enter a company identifier to find similar deals")

        identifier_type = st.selectbox(
            "Select an identifier type",
            ["website_url", "website_domain", "linkedin_url", "crunchbase_url", "pitchbook_url",
            "twitter_url", "instagram_url", "facebook_url", "angellist_url",
            "monster_url", "indeed_url", "stackoverflow_url"]
        )
        identifier_value = st.text_input("Enter the identifier value (e.g., domain or URL)")

        if st.button("Enrich Company & Find Similar Deals", disabled=(not identifier_value.strip())):
            try:
                st.info("Enriching company...")
                main_company = enrich_company(identifier_type, identifier_value.strip())

                st.success(f"Enriched: {main_company.get('legal_name') or main_company.get('name')}")

                urn = main_company.get("entity_urn")
                if urn:
                    st.info("Fetching similar companies...")
                    similar_urns = get_similar_companies(urn)
                    similar_companies = batch_get_companies(similar_urns)

                    st.subheader("Main Company")
                    st.markdown(f"**Name:** {main_company.get('legal_name') or main_company.get('name')}")
                    st.markdown(f"**Description:** {main_company.get('description','No description available.')}")
                    website = main_company.get('website', {}).get('url', '')
                    if website:
                        st.markdown(f"**Website:** [{website}]({website})")
                    
                    stage = main_company.get('stage')
                    if stage:
                        st.markdown(f"**Stage:** {stage}")

                    ownership = main_company.get('ownership_status')
                    if ownership:
                        st.markdown(f"**Ownership:** {ownership}")

                    location = main_company.get('location', {})
                    city = location.get('city', '')
                    region = location.get('region', '')
                    country = location.get('country', '')
                    if any([city, region, country]):
                        st.markdown(f"**Location:** {city} {region} {country}")

                    st.markdown("---")
                    st.subheader("Similar Companies")
                    for company in similar_companies:
                        st.markdown(f"### {company.get('legal_name') or company.get('name')}")
                        st.markdown(f"**Description:** {company.get('description','No description available.')}")
                        
                        website_info = company.get('website', {})
                        if website_info:
                            url = website_info.get('url') or ""
                            domain = website_info.get('domain') or ""
                            st.markdown(f"**Website:** [{domain}]({url})" if url else f"**Website:** {domain}")

                        stage = company.get('stage')
                        if stage:
                            st.markdown(f"**Stage:** {stage}")

                        ownership = company.get('ownership_status')
                        if ownership:
                            st.markdown(f"**Ownership:** {ownership}")

                        location = company.get('location', {})
                        city = location.get('city', '')
                        region = location.get('region', '')
                        country = location.get('country', '')
                        if any([city, region, country]):
                            st.markdown(f"**Location:** {city} {region} {country}")

                        st.markdown("---")

                    # PDF download
                    pdf_buffer = generate_companies_report(main_company, similar_companies)
                    with st.expander("📄 Download PDF Report"):
                        st.download_button(
                            label="📥 Download Main Company & Similar Deals Report",
                            data=pdf_buffer,
                            file_name=f"{main_company.get('name','company')}_similar_companies.pdf",
                            mime="application/pdf"
                        )

                else:
                    st.error("Could not retrieve the company's URN for similar company lookup.")

            except Exception as e:
                st.error(f"Error: {e}")


# --- Main Application Logic ---
st.set_page_config(
    page_title="Y+ Ventures Portal", layout="wide",
    page_icon="icons/yplusventures_logo.jpeg"
)

st.sidebar.title("Y+ Ventures Portal")
app_mode = st.sidebar.radio("Select Interface", ["Founder Interface", "VC Interface"])

VC_PASSWORD = st.secrets["passwords"]["vc_password"]

if app_mode == "Founder Interface":
    from gemini_new_app56 import founder_form
    founder_form()

else:  # VC Interface
    # Check if user is authenticated already
    if "vc_authenticated" not in st.session_state:
        st.session_state.vc_authenticated = False

    if not st.session_state.vc_authenticated:
        st.title("🔒 VC Access Restricted")
        password_input = st.text_input("Enter VC Password:", type="password")

        if st.button("Login"):
            if password_input == VC_PASSWORD:
                st.session_state.vc_authenticated = True
                st.success("Access granted ✅")
                st.rerun()  # reload to show VC interface
            else:
                st.error("Incorrect password. Try again.")

    else:
        vc_main_app_logic()
